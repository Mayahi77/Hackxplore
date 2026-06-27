import uuid
import re
import io
import os
import shutil
import subprocess
import tempfile
from pathlib import Path


# ── Cross-discipline knowledge schema ──────────────────────────────────────
# Tags let the knowledge base scale beyond two-stroke engines. Every indexed
# segment carries a `discipline` so the same tool can serve mechanical,
# electrical, combustion/physics and simulation knowledge side by side.
DISCIPLINES = [
    "Mechanical",
    "Electrical",
    "Physics/Combustion",
    "Simulation Data",
    "General",
]

_DISCIPLINE_KEYWORDS: dict[str, list[str]] = {
    "Mechanical": [
        "piston", "crankshaft", "bearing", "gear", "torque", "cylinder",
        "valve", "gasket", "seal", "wear", "vibration", "shaft", "housing",
        "mount", "fatigue", "stress", "clearance", "lubrication", "friction",
    ],
    "Electrical": [
        "ignition", "spark", "coil", "stator", "voltage", "current", "wiring",
        "sensor", "ecu", "battery", "circuit", "capacitor", "cdi", "magneto",
        "resistance", "ground", "relay", "harness",
    ],
    "Physics/Combustion": [
        "combustion", "detonation", "temperature", "thermodynam", "pressure",
        "fuel", "air ratio", "mixture", "flame", "heat", "knock", "octane",
        "stoichio", "expansion", "exhaust gas", "compression",
    ],
    "Simulation Data": [
        "simulation", "cfd", "fea", "mesh", "solver", "ansys", "matlab",
        "dataset", "numerical", "boundary condition", "fem", "model run",
        "convergence", "residual",
    ],
}


def classify_discipline(text: str) -> str:
    """Best-effort keyword classifier. Returns 'General' when nothing matches."""
    t = text.lower()
    best, best_score = "General", 0
    for discipline, keywords in _DISCIPLINE_KEYWORDS.items():
        score = sum(t.count(kw) for kw in keywords)
        if score > best_score:
            best, best_score = discipline, score
    return best


def normalize_discipline(value: str | None) -> str | None:
    """Map a user-supplied discipline to a known value, else None (auto)."""
    if not value:
        return None
    value = value.strip()
    for d in DISCIPLINES:
        if d.lower() == value.lower():
            return d
    return None


# ── Chunking helpers ───────────────────────────────────────────────────────
def _chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> list[str]:
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        if chunk.strip():
            chunks.append(chunk.strip())
        start += chunk_size - overlap
    return chunks


def _clean(text: str) -> str:
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def build_text_chunks(
    text: str,
    filename: str,
    doc_type: str,
    discipline: str | None = None,
    extra_meta: dict | None = None,
) -> list[dict]:
    """Reusable builder: clean text -> chunk -> stamped records.

    Used for free text such as transcribed images and typed/voice notes.
    """
    text = _clean(text)
    resolved = discipline or classify_discipline(text)
    chunks: list[dict] = []
    for i, chunk in enumerate(_chunk_text(text)):
        meta = {
            "source": filename,
            "filename": filename,
            "page": 1,
            "chunk_index": i,
            "doc_type": doc_type,
            "discipline": resolved,
        }
        if extra_meta:
            meta.update(extra_meta)
        chunks.append({"id": str(uuid.uuid4()), "text": chunk, "metadata": meta})
    return chunks


# ── Vision helpers for presentations ──────────────────────────────────────
def _pptx_vision_enabled() -> bool:
    flag = os.getenv("PPTX_VISION_ENABLED", "true").strip().lower()
    if flag in {"0", "false", "no", "off"}:
        return False
    key = os.getenv("GOOGLE_API_KEY", "").strip()
    return bool(key and key != "your_google_api_key_here")


def _pptx_vision_max_items() -> int:
    try:
        return max(1, int(os.getenv("PPTX_VISION_MAX_SLIDES", "20")))
    except ValueError:
        return 20


def _find_soffice() -> str | None:
    configured = os.getenv("LIBREOFFICE_BIN", "").strip()
    candidates = [
        configured,
        shutil.which("soffice") or "",
        shutil.which("libreoffice") or "",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]
    for candidate in candidates:
        if candidate and os.path.isfile(candidate) and os.access(candidate, os.X_OK):
            return candidate
    return None


def _render_pptx_slides(file_content: bytes) -> list[tuple[int, bytes]]:
    """Render PPTX slides to PNG bytes when LibreOffice is available."""
    soffice = _find_soffice()
    if not soffice:
        return []

    timeout = int(os.getenv("PPTX_RENDER_TIMEOUT_SECONDS", "90"))
    with tempfile.TemporaryDirectory(prefix="pptx-render-") as tmp:
        tmp_path = Path(tmp)
        pptx_path = tmp_path / "deck.pptx"
        pptx_path.write_bytes(file_content)

        subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(tmp_path),
                str(pptx_path),
            ],
            cwd=tmp_path,
            check=True,
            capture_output=True,
            timeout=timeout,
        )

        pdfs = sorted(tmp_path.glob("*.pdf"))
        if not pdfs:
            return []

        import fitz  # pymupdf

        rendered: list[tuple[int, bytes]] = []
        doc = fitz.open(pdfs[0])
        zoom = float(os.getenv("PPTX_RENDER_ZOOM", "2.0"))
        matrix = fitz.Matrix(zoom, zoom)
        for page_num, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            rendered.append((page_num, pix.tobytes("png")))
        return rendered


def _image_bytes_for_vision(image) -> tuple[bytes, str] | None:
    """Return image bytes and a Gemini-supported mime type."""
    mime = (getattr(image, "content_type", "") or "").lower()
    blob = image.blob
    if mime in {"image/png", "image/jpeg", "image/webp", "image/gif", "image/bmp"}:
        return blob, mime

    try:
        from PIL import Image

        with Image.open(io.BytesIO(blob)) as im:
            out = io.BytesIO()
            im.convert("RGB").save(out, format="PNG")
            return out.getvalue(), "image/png"
    except Exception:
        return None


def _embedded_pptx_images(file_content: bytes) -> list[tuple[int, bytes, str, str]]:
    """Extract embedded raster images from slides as a fallback to full rendering."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_content))
    items: list[tuple[int, bytes, str, str]] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        image_index = 0
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) != 13 or not hasattr(shape, "image"):
                continue
            converted = _image_bytes_for_vision(shape.image)
            if not converted:
                continue
            image_index += 1
            image_bytes, mime = converted
            items.append((slide_num, image_bytes, mime, f"embedded image {image_index}"))
    return items


def _pptx_visual_chunks(file_content: bytes, filename: str) -> list[dict]:
    """Best-effort slide visual understanding for diagrams/screenshots.

    Full slide rendering catches drawn arrows/shapes/layout when LibreOffice is
    installed. Embedded-image extraction catches screenshots and raster diagrams
    when rendering is unavailable.
    """
    if not _pptx_vision_enabled():
        return []

    from app.services.vision import extract_image_knowledge

    visual_items: list[tuple[int, bytes, str, str]] = []
    extraction_mode = "pptx_slide_render"
    try:
        rendered = _render_pptx_slides(file_content)
        visual_items = [
            (slide_num, image_bytes, "image/png", "rendered full slide")
            for slide_num, image_bytes in rendered
        ]
    except Exception:
        visual_items = []

    if not visual_items:
        extraction_mode = "pptx_embedded_image"
        visual_items = _embedded_pptx_images(file_content)

    chunks: list[dict] = []
    for visual_index, (slide_num, image_bytes, mime, label) in enumerate(
        visual_items[:_pptx_vision_max_items()],
        start=1,
    ):
        hint = (
            f"PowerPoint file '{filename}', slide {slide_num}, {label}. "
            "Focus on engineering diagrams, arrows, labels, dimensions, plots, tables, "
            "component relationships and any visible technical values."
        )
        try:
            extracted = extract_image_knowledge(image_bytes, mime, hint=hint)
        except Exception:
            continue
        if not extracted:
            continue

        text = _clean(
            f"PowerPoint visual content from {filename}, slide {slide_num}, {label}. "
            f"{extracted}"
        )
        for i, chunk in enumerate(_chunk_text(text, chunk_size=500, overlap=60)):
            chunks.append({
                "id": str(uuid.uuid4()),
                "text": chunk,
                "metadata": {
                    "source": filename,
                    "filename": filename,
                    "page": slide_num,
                    "chunk_index": 10000 + visual_index * 100 + i,
                    "doc_type": "pptx",
                    "modality": "visual",
                    "extraction": extraction_mode,
                },
            })
    return chunks


# ── Per-format parsers ─────────────────────────────────────────────────────
def process_pdf(file_content: bytes, filename: str) -> list[dict]:
    import fitz  # pymupdf

    chunks = []
    doc = fitz.open(stream=file_content, filetype="pdf")
    for page_num, page in enumerate(doc, start=1):
        text = _clean(page.get_text())
        if len(text) < 50:
            continue
        for i, chunk in enumerate(_chunk_text(text)):
            chunks.append({
                "id": str(uuid.uuid4()),
                "text": chunk,
                "metadata": {
                    "source": filename,
                    "filename": filename,
                    "page": page_num,
                    "chunk_index": i,
                    "doc_type": "pdf",
                },
            })
    return chunks


def process_docx(file_content: bytes, filename: str) -> list[dict]:
    from docx import Document

    doc = Document(io.BytesIO(file_content))
    full_text = _clean(" ".join(p.text for p in doc.paragraphs if p.text.strip()))
    chunks = []
    for i, chunk in enumerate(_chunk_text(full_text)):
        chunks.append({
            "id": str(uuid.uuid4()),
            "text": chunk,
            "metadata": {
                "source": filename,
                "filename": filename,
                "page": 1,
                "chunk_index": i,
                "doc_type": "docx",
            },
        })
    return chunks


def process_pptx(file_content: bytes, filename: str) -> list[dict]:
    """Parse native PowerPoint files directly — text, tables and speaker notes.

    Much internal engineering knowledge lives in presentations rather than PDFs,
    so we index each slide individually (page == slide number).
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_content))
    chunks: list[dict] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    row_text = " | ".join(c for c in cells if c)
                    if row_text:
                        parts.append(row_text)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Speaker notes: {notes}")

        text = _clean(" ".join(parts))
        if len(text) < 10:
            continue
        for i, chunk in enumerate(_chunk_text(text)):
            chunks.append({
                "id": str(uuid.uuid4()),
                "text": chunk,
                "metadata": {
                    "source": filename,
                    "filename": filename,
                    "page": slide_num,
                    "chunk_index": i,
                    "doc_type": "pptx",
                    "modality": "text",
                },
            })
    chunks.extend(_pptx_visual_chunks(file_content, filename))
    return chunks


def process_xlsx(file_content: bytes, filename: str) -> list[dict]:
    import pandas as pd

    chunks = []
    xls = pd.ExcelFile(io.BytesIO(file_content))
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name).fillna("")
        text = _clean(f"Sheet: {sheet_name}\n" + df.to_string(index=False))
        for i, chunk in enumerate(_chunk_text(text)):
            chunks.append({
                "id": str(uuid.uuid4()),
                "text": chunk,
                "metadata": {
                    "source": filename,
                    "filename": filename,
                    "page": 1,
                    "chunk_index": i,
                    "doc_type": "xlsx",
                    "sheet": sheet_name,
                },
            })
    return chunks


def process_txt(file_content: bytes, filename: str) -> list[dict]:
    text = _clean(file_content.decode("utf-8", errors="ignore"))
    chunks = []
    for i, chunk in enumerate(_chunk_text(text)):
        chunks.append({
            "id": str(uuid.uuid4()),
            "text": chunk,
            "metadata": {
                "source": filename,
                "filename": filename,
                "page": 1,
                "chunk_index": i,
                "doc_type": "txt",
            },
        })
    return chunks


def process_file(
    file_content: bytes,
    filename: str,
    discipline: str | None = None,
) -> list[dict]:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        chunks = process_pdf(file_content, filename)
    elif ext == ".docx":
        chunks = process_docx(file_content, filename)
    elif ext == ".pptx":
        chunks = process_pptx(file_content, filename)
    elif ext == ".xlsx":
        chunks = process_xlsx(file_content, filename)
    elif ext in (".txt", ".md"):
        chunks = process_txt(file_content, filename)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    # Stamp a discipline on every segment — explicit override or auto-classified.
    resolved = normalize_discipline(discipline)
    for c in chunks:
        c["metadata"]["discipline"] = resolved or classify_discipline(c["text"])
    return chunks
